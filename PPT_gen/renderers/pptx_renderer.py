from pptx import Presentation
from pptx.util import Inches


def render_pptx(deck, output_path, visuals=None):
    prs = Presentation()

    for i, slide_data in enumerate(deck.slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data.title
        slide.placeholders[1].text = "\n".join(slide_data.bullets)

        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes

        if visuals and i < len(visuals):
            slide.shapes.add_picture(
                str(visuals[i]),
                Inches(5.5),
                Inches(1.5),
                width=Inches(4)
            )

    prs.save(output_path)